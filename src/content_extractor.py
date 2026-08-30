# -*- coding: utf-8 -*-
"""
文档内容提取器 — 从各种文件格式中提取纯文本内容。

支持格式：
- PDF (.pdf) — 使用 pypdf
- Word (.docx) — 使用 python-docx
- Excel (.xlsx) — 使用 openpyxl
- PowerPoint (.pptx) — 使用 python-pptx
- 纯文本 (.txt, .md, .csv, .log, .json, .xml, .html)
- RTF (.rtf) — 使用 striprtf

所有提取器返回纯文本字符串，用于后续索引和预览。
"""
from __future__ import annotations

import os
import re
from pathlib import Path
from typing import Optional, Dict, List


# ─── 支持的文件类型映射 ──────────────────────────────────────────────────────
SUPPORTED_EXTENSIONS = {
    # 文档
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".rtf": "rtf",
    # 纯文本
    ".txt": "text",
    ".md": "text",
    ".csv": "text",
    ".log": "text",
    ".json": "text",
    ".xml": "text",
    ".html": "text",
    ".htm": "text",
    ".ini": "text",
    ".cfg": "text",
    ".conf": "text",
    ".yaml": "text",
    ".yml": "text",
    ".py": "text",
    ".js": "text",
    ".ts": "text",
    ".java": "text",
    ".c": "text",
    ".cpp": "text",
    ".h": "text",
    ".cs": "text",
    ".go": "text",
    ".rb": "text",
    ".php": "text",
    ".sh": "text",
    ".bat": "text",
    ".ps1": "text",
    ".sql": "text",
    ".r": "text",
    ".m": "text",
}


def _read_text_file(path: str, max_chars: int = 1_000_000) -> str:
    """读取纯文本文件，自动检测编码。"""
    encodings = ["utf-8", "gbk", "gb2312", "gb18030", "big5", "latin-1"]
    for enc in encodings:
        try:
            with open(path, "r", encoding=enc, errors="strict") as f:
                content = f.read(max_chars)
                return content
        except (UnicodeDecodeError, UnicodeError):
            continue
        except Exception:
            return ""
    # 最后尝试用 errors='replace'
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(max_chars)
    except Exception:
        return ""


def _extract_pdf(path: str, max_chars: int = 1_000_000) -> str:
    """从 PDF 提取文本。"""
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                return ""
        pages = []
        total_chars = 0
        for page in reader.pages:
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                pages.append(text)
                total_chars += len(text)
                if total_chars >= max_chars:
                    break
        return "\n\n".join(pages)
    except Exception:
        return ""


def _extract_docx(path: str, max_chars: int = 1_000_000) -> str:
    """从 DOCX 提取文本（段落 + 表格）。"""
    try:
        from docx import Document
        doc = Document(path)
        parts = []
        total_chars = 0

        # 提取段落
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
                total_chars += len(text)
                if total_chars >= max_chars:
                    break

        # 提取表格
        if total_chars < max_chars:
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)
                        total_chars += len(row_text)
                        if total_chars >= max_chars:
                            break

        return "\n".join(parts)
    except Exception:
        return ""


def _extract_xlsx(path: str, max_chars: int = 1_000_000) -> str:
    """从 XLSX 提取文本（所有 sheet 的所有单元格）。"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts = []
        total_chars = 0

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_parts = [f"[Sheet: {sheet_name}]"]
            for row in ws.iter_rows(values_only=True):
                row_values = [str(cell) for cell in row if cell is not None]
                if row_values:
                    row_text = " | ".join(row_values)
                    sheet_parts.append(row_text)
                    total_chars += len(row_text)
                    if total_chars >= max_chars:
                        break
            if len(sheet_parts) > 1:
                parts.extend(sheet_parts)
            if total_chars >= max_chars:
                break

        wb.close()
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_pptx(path: str, max_chars: int = 1_000_000) -> str:
    """从 PPTX 提取文本（所有幻灯片的所有文本框）。"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        total_chars = 0

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {i}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)
                            total_chars += len(text)
                            if total_chars >= max_chars:
                                break
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_parts.append(row_text)
                            total_chars += len(row_text)
                if total_chars >= max_chars:
                    break
            if len(slide_parts) > 1:
                parts.extend(slide_parts)
            if total_chars >= max_chars:
                break

        return "\n".join(parts)
    except Exception:
        return ""


def _extract_rtf(path: str, max_chars: int = 1_000_000) -> str:
    """从 RTF 提取文本。"""
    try:
        from striprtf.striprtf import rtf_to_text
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            rtf_content = f.read(max_chars * 2)
        text = rtf_to_text(rtf_content)
        return text[:max_chars]
    except ImportError:
        # striprtf 未安装，尝试简单提取
        return _extract_rtf_simple(path, max_chars)
    except Exception:
        return ""


def _extract_rtf_simple(path: str, max_chars: int = 1_000_000) -> str:
    """简单的 RTF 文本提取（不依赖 striprtf）。"""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read(max_chars * 2)
        # 移除 RTF 控制字
        text = re.sub(r"\\[a-z]+\d*\s?", "", content)
        # 移除大括号
        text = re.sub(r"[{}]", "", text)
        # 移除多余空白
        text = re.sub(r"\s+", " ", text).strip()
        return text[:max_chars]
    except Exception:
        return ""


# ─── 统一提取接口 ─────────────────────────────────────────────────────────────
class ContentExtractor:
    """
    文档内容提取器。
    根据文件扩展名选择合适的提取方法。
    """

    def __init__(self):
        self._extractors = {
            "pdf": _extract_pdf,
            "docx": _extract_docx,
            "xlsx": _extract_xlsx,
            "pptx": _extract_pptx,
            "rtf": _extract_rtf,
            "text": _read_text_file,
        }

    def extract(self, path: str, max_chars: int = 1_000_000) -> str:
        """
        从文件中提取文本内容。

        Args:
            path: 文件路径
            max_chars: 最大提取字符数

        Returns:
            提取的文本内容，失败返回空字符串
        """
        if not os.path.isfile(path):
            return ""

        ext = os.path.splitext(path)[1].lower()
        file_type = SUPPORTED_EXTENSIONS.get(ext)
        if not file_type:
            return ""

        extractor = self._extractors.get(file_type)
        if not extractor:
            return ""

        try:
            return extractor(path, max_chars)
        except Exception:
            return ""

    def can_extract(self, path: str) -> bool:
        """检查是否支持提取该文件。"""
        ext = os.path.splitext(path)[1].lower()
        return ext in SUPPORTED_EXTENSIONS

    def get_supported_extensions(self) -> List[str]:
        """返回支持的扩展名列表。"""
        return sorted(SUPPORTED_EXTENSIONS.keys())

    def extract_preview(self, path: str, max_chars: int = 5000) -> str:
        """提取文件预览内容（较短）。"""
        return self.extract(path, max_chars)
